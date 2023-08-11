from django.contrib.auth import authenticate
from django.shortcuts import render, redirect
from django.contrib.auth.models import User, auth
from django.contrib.auth.decorators import login_required
from .models import *
from pptx import Presentation
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from .models import *
import openai
from django.core.files.storage import FileSystemStorage
from pptx import Presentation
import os
from django.conf import settings



# Create your views here.

def home(request):
    template = 'home.html'
    
    context = {}
        
    return render(request, template, context) # Return for Templates


def dashboard(request):
    template = 'dashboard.html'
    chat_all = Chat.objects.all()
    context = {'chat_all':chat_all}
    return render(request, template, context) # Return for Templates



@login_required(login_url='login')
def new_chat(request):
    chat_all = Chat.objects.all()
    context = {'chat_all':chat_all}
    if request.method == 'POST':
        title = request.POST.get('title')
        pptx_file = request.FILES['pptx_file']
        fs = FileSystemStorage()
        filename = fs.save(pptx_file.name, pptx_file)
        pptx_path = os.path.join(settings.MEDIA_ROOT, filename)

        presentation = Presentation(pptx_path)
        chat_cookie = ""

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chat_cookie += shape.text + "\n"
                    
        print(chat_cookie)
        chat = Chat.objects.create(title=title, chat_cookie=chat_cookie, user=request.user)
        return redirect('chat_detail', chat_id=chat.id)
    return render(request, 'new_chat.html', context)


@login_required(login_url='login')
def cookie_data(request, chat_id):
    chat = get_object_or_404(Chat, id=chat_id)
    ChatData = ChatMessage.objects.filter(chat=chat)
    chat_all = Chat.objects.all()
    content = chat.chat_cookie.replace('—','</br></br>')
    context = {'chat_all':chat_all, 'chat': chat, 'ChatData': ChatData, "content":content}
    return render(request, 'cookie.html', context) 


@login_required(login_url='login')
def chat_detail(request, chat_id):
    chat = get_object_or_404(Chat, id=chat_id)
    ChatData = ChatMessage.objects.filter(chat=chat)
    
    openai_key = OpenAIKey.objects.all()[0].key
    print(openai_key)
    openai.api_key = openai_key
    
    chat_all = Chat.objects.all()
    context = {'chat_all':chat_all, 'chat': chat, 'ChatData': ChatData}
    
    if request.method == 'POST':
        prompt = request.POST.get('prompt')
        sender = request.user
        
        openai.api_key = openai_key
        response = openai.ChatCompletion.create(model='gpt-3.5-turbo', messages=[{"role": "system", "content": chat.chat_cookie},{"role": "user", "content": prompt}])
        result = ''
        for choice in response.choices:
            result += choice.message.content
    
        replay = result
        ChatMessage.objects.create(chat=chat, sender=sender, prompt=prompt, replay=replay)
        
        return redirect('chat_detail', chat_id=chat.id)

    return render(request, 'chat_detail.html', context)


@login_required(login_url='login/')  # login/  is custom login URL path
def delete_chat(request, chat_id):   # ```data_id``` should be pass in url as <data_id>
    data = Chat.objects.get(pk=chat_id)
    data.delete()
    return redirect('/dashboard')